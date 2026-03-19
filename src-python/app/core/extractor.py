"""
Text extraction from various file types.

Each extractor handles a set of file extensions and returns the raw text
content. The registry pattern makes it easy to add new extractors.
"""

from __future__ import annotations

from pathlib import Path

import structlog

logger = structlog.get_logger()


class ExtractionError(Exception):
    """Raised when text extraction fails."""


def extract_text(file_path: Path) -> str:
    """
    Extract text content from a file based on its extension.

    Returns the extracted text, or raises ExtractionError if extraction fails.
    """
    ext = file_path.suffix.lower()

    extractors = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".doc": _extract_docx,     # python-docx handles .doc too
        ".pptx": _extract_pptx,
        ".xlsx": _extract_xlsx,
        ".csv": _extract_csv,
        ".tsv": _extract_csv,
    }

    extractor = extractors.get(ext, _extract_plaintext)

    try:
        text = extractor(file_path)
        return text.strip()
    except Exception as e:
        raise ExtractionError(f"Failed to extract {file_path}: {e}") from e


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF using PyMuPDF (fitz)."""
    import fitz

    doc = fitz.open(str(path))
    pages = []
    for page in doc:
        text = page.get_text("text")
        if text.strip():
            pages.append(text)
    doc.close()

    if not pages:
        logger.debug("extraction.pdf.empty", path=str(path))
        return ""

    return "\n\n".join(pages)


def _extract_docx(path: Path) -> str:
    """Extract text from Word documents."""
    from docx import Document

    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_pptx(path: Path) -> str:
    """Extract text from PowerPoint presentations."""
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append("\n".join(texts))

    return "\n\n---\n\n".join(slides)


def _extract_xlsx(path: Path) -> str:
    """Extract text from Excel spreadsheets."""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"[Sheet: {ws.title}]\n" + "\n".join(rows))
    wb.close()

    return "\n\n".join(sheets)


def _extract_csv(path: Path) -> str:
    """Extract text from CSV/TSV files."""
    import csv
    import chardet

    # Detect encoding
    raw = path.read_bytes()[:10000]
    detected = chardet.detect(raw)
    encoding = detected.get("encoding", "utf-8") or "utf-8"

    delimiter = "\t" if path.suffix.lower() == ".tsv" else ","

    rows = []
    with open(path, "r", encoding=encoding, errors="replace") as f:
        reader = csv.reader(f, delimiter=delimiter)
        for row in reader:
            cells = [c.strip() for c in row if c.strip()]
            if cells:
                rows.append(" | ".join(cells))

    return "\n".join(rows)


def _extract_plaintext(path: Path) -> str:
    """
    Extract text from plaintext files (code, markdown, config, etc.).

    Handles encoding detection for non-UTF-8 files.
    """
    import chardet

    raw = path.read_bytes()

    # Quick check: try UTF-8 first (most common)
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        pass

    # Fall back to detected encoding
    detected = chardet.detect(raw[:10000])
    encoding = detected.get("encoding", "latin-1") or "latin-1"

    return raw.decode(encoding, errors="replace")
